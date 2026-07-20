"""
Multi-format document parser.
Supports: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx)
"""
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


def _extract_page_text(page) -> str:
    """
    Extract page text, adapting to PDFs with very tight glyph spacing where the
    default tolerance merges words together ("Understandingthesefundamentals").
    If the space ratio looks abnormally low, retry with a tighter x_tolerance.
    """
    text = page.extract_text() or ""
    if text and (text.count(" ") / max(len(text), 1)) < 0.05:
        retry = page.extract_text(x_tolerance=1.5) or ""
        if retry.count(" ") > text.count(" "):
            return retry
    return text


def parse_pdf(file_path: str) -> str:
    import pdfplumber
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = _extract_page_text(page)
            if page_text:
                text += f"\n[Page {i+1}]\n{page_text}\n"
            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    clean = [str(c).strip() if c else "" for c in row]
                    text += " | ".join(clean) + "\n"
    return text.strip()


def parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def parse_xlsx(file_path: str) -> str:
    import pandas as pd
    text = ""
    xl = pd.ExcelFile(file_path)
    for sheet in xl.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet, dtype=str).fillna("")
        text += f"\n--- Sheet: {sheet} ---\n"
        # Header
        text += " | ".join(str(c) for c in df.columns) + "\n"
        # Rows
        for _, row in df.iterrows():
            row_str = " | ".join(str(v) for v in row.values)
            if row_str.strip(" |"):
                text += row_str + "\n"
    return text.strip()


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"\n[Slide {i+1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def parse_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ── Dispatcher ─────────────────────────────────────────────────────────────────
PARSERS = {
    ".pdf":  parse_pdf,
    ".docx": parse_docx,
    ".doc":  parse_docx,
    ".xlsx": parse_xlsx,
    ".xls":  parse_xlsx,
    ".pptx": parse_pptx,
    ".ppt":  parse_pptx,
    ".txt":  parse_txt,
}

SUPPORTED_FORMATS = list(PARSERS.keys())


def parse_document(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"Unsupported format '{ext}'. Supported: {SUPPORTED_FORMATS}")
    logger.info(f"Parsing {Path(file_path).name} as {ext}")
    return parser(file_path)
